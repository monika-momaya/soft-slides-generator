"""
pptx_compositor.py

Generates an EDITABLE PowerPoint (.pptx) slide for a session, using the
same layout engine and theme detection as the PNG compositor, but placing
real PowerPoint text boxes and picture shapes instead of drawing pixels.

Why this exists (vs. slide_compositor.py): event teams need to make small
on-the-spot edits during a live event (a speaker drops out, a title typo,
etc.) - a flat PNG can't be edited in PowerPoint, but a deck with real
text boxes and picture placeholders can.

Design decisions carried over from earlier discussion:
- The branded template PNG becomes the SLIDE BACKGROUND IMAGE (not
  recreated as native shapes) - designers own branding, we own data.
- Speaker photos keep their exact custom mask silhouette (circle, hexagon,
  whatever the designer drew) by pre-flattening the crop with Pillow
  (reusing mask_parser + the same masking logic as the PNG path), then
  inserting that flattened image as a PowerPoint picture. This trades
  "PowerPoint can natively re-crop this shape" for "the shape is pixel-
  perfect to what the designer intended" - a deliberate tradeoff per
  project decision.
- Text (session title, speaker name/title/company, role labels) becomes
  REAL editable PowerPoint text boxes, with font/size/color/bold matching
  the same rules as the PNG version (PPT-equivalent pt sizes, theme-aware
  coloring, bold names, etc.) - these can be edited live in PowerPoint.

One slide = one session. Multi-session decks are a planned follow-up
(see project notes) - for now, build_pptx_slide() adds one slide to a
given Presentation object, so multi-session support later is just calling
this function in a loop against the same Presentation.
"""

import io
from typing import List
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from core.layout_engine import Slot
from core.mask_parser import SlotShape
from core.theme_detector import detect_theme, TextTheme
from core.slide_compositor import (
    _apply_mask_shape, SESSION_NAME_IDEAL_PT, SESSION_NAME_MIN_PT, META_LINE_PT,
    SPEAKER_NAME_PT, TITLE_COMPANY_PT, ROLE_LABEL_PT, DEFAULT_LABEL_VALUES_TO_HIDE,
    CANVAS_TOP_RATIO,
)

# Standard PowerPoint 16:9 widescreen slide dimensions
SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5


def _rgb(color_tuple) -> RGBColor:
    return RGBColor(color_tuple[0], color_tuple[1], color_tuple[2])


def _save_pil_image_to_stream(img: Image.Image, preserve_alpha: bool = False) -> io.BytesIO:
    """
    Saves a PIL image to an in-memory PNG stream for embedding in the
    PPTX. preserve_alpha=True keeps transparency (required for masked
    speaker photos, where transparent areas outside the mask shape must
    stay transparent - not flattened to black/white) - False flattens to
    RGB (fine for the opaque background template image).
    """
    stream = io.BytesIO()
    if preserve_alpha:
        img.convert("RGBA").save(stream, format="PNG")
    else:
        img.convert("RGB").save(stream, format="PNG")
    stream.seek(0)
    return stream


def new_presentation() -> Presentation:
    """Create a blank Presentation sized to standard 16:9 widescreen."""
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_WIDTH_IN * 914400))
    prs.slide_height = Emu(int(SLIDE_HEIGHT_IN * 914400))
    return prs


def _add_text_box(slide, left_in, top_in, width_in, height_in, text, font_pt,
                   bold, color_tuple, align=PP_ALIGN.CENTER, font_name="Calibri",
                   anchor=MSO_ANCHOR.TOP, autosize_shrink=True):
    """
    Adds a text box with no fill/border, sized to (width_in, height_in),
    with the given single-line/paragraph text. Word wrap is OFF for
    speaker names per project decision (shrink-to-fit, never wrap) - the
    caller is responsible for not feeding text wider than the box; for
    safety we still set word_wrap=True at the box level (PowerPoint text
    boxes need this to avoid clipping unexpectedly long strings the user
    later edits in), but font sizing is chosen to fit on one line at
    render time.
    """
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = _rgb(color_tuple)
    return box


def _measure_text_width_in(text: str, font_pt: float, bold: bool) -> float:
    """
    Rough text width estimate in inches, used to decide if a speaker
    name needs a smaller font to fit its slot on one line (PowerPoint
    doesn't expose text measurement without rendering, so this uses a
    standard average-character-width heuristic - generous enough to be
    safe, not pixel-exact, which is fine since PowerPoint text boxes
    reflow live anyway if a user edits the text afterward).
    """
    avg_char_width_ratio = 0.56 if bold else 0.52
    return len(text) * font_pt * avg_char_width_ratio / 72.0


def _fit_font_pt(text: str, max_width_in: float, ideal_pt: float, min_pt: float, bold: bool) -> float:
    """Shrink font size (never below min_pt) until text fits max_width_in on one line."""
    size = ideal_pt
    while size > min_pt:
        if _measure_text_width_in(text, size, bold) <= max_width_in:
            return size
        size -= 0.5
    return min_pt


def _shared_name_font_pt(names: List[str], max_width_in: float, ideal_pt: float, min_pt: float) -> float:
    """
    Computes ONE font size shared by ALL speaker names on the slide, rather
    than fitting each name independently. This matters for real-world
    PowerPoint editing: if every name starts at the same size, selecting
    all the name text boxes and changing the font size in PowerPoint's
    toolbar actually behaves the way someone expects ("change this one
    number, they all update") - independent per-name sizing made that
    impossible, since boxes could already differ before any manual edit.

    The shared size is the smallest size needed to fit the LONGEST name
    (bold) within max_width_in - i.e. driven by whichever name is hardest
    to fit, same as the previous per-name logic would have picked for that
    one name, but now applied uniformly to every name on the slide.
    """
    if not names:
        return ideal_pt
    size = ideal_pt
    while size > min_pt:
        if all(_measure_text_width_in(n, size, bold=True) <= max_width_in for n in names):
            return size
        size -= 0.5
    return min_pt


def _wrap_text_simple(text: str, max_width_in: float, font_pt: float, bold: bool, max_lines: int = 2) -> str:
    """
    Simple width-based word wrap (greedy line-fill, breaks wherever a line
    hits max_width_in - NOT phrase-aware, per project decision). Returns
    the text with '\\n' inserted at break points, capped at max_lines
    (remaining words are appended to the last line rather than truncated,
    so no content is ever silently dropped).
    """
    words = text.split()
    if not words:
        return text

    lines = []
    current = words[0]
    for word in words[1:]:
        candidate = f"{current} {word}"
        if _measure_text_width_in(candidate, font_pt, bold) <= max_width_in:
            current = candidate
        else:
            lines.append(current)
            current = word
        if len(lines) == max_lines - 1:
            # Last allowed line - dump everything remaining onto it rather
            # than dropping words, even if it overflows slightly; PowerPoint
            # text boxes can be manually widened/edited if this looks tight.
            remaining_idx = words.index(word) + 1
            current = " ".join([current] + words[remaining_idx:])
            break
    lines.append(current)
    return "\n".join(lines[:max_lines])


def _add_multiline_text_box(slide, left_in, top_in, width_in, height_in, text, font_pt,
                             bold, color_tuple, align=PP_ALIGN.CENTER, font_name="Calibri",
                             anchor=MSO_ANCHOR.TOP):
    """Like _add_text_box, but splits on '\\n' into separate paragraphs (for wrapped text)."""
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_pt)
        run.font.bold = bold
        run.font.name = font_name
        run.font.color.rgb = _rgb(color_tuple)
    return box


def _build_title_company_lines(title: str, company: str, max_width_in: float, font_pt: float) -> str:
    """
    Decides how to lay out the title/company caption:
    - If both are present and fit together on one line as "Title, Company",
      combine them that way (e.g. "Founder, Infosys") - matches the
      "smaller names" example from project feedback.
    - Otherwise, title and company are each wrapped independently (simple
      width-based wrap, 2 lines max each) and stacked.
    Returns a '\\n'-joined string ready for _add_multiline_text_box.
    """
    if title and company:
        combined = f"{title}, {company}"
        if _measure_text_width_in(combined, font_pt, bold=False) <= max_width_in:
            return combined
        # Doesn't fit combined - wrap each piece independently and stack
        wrapped_title = _wrap_text_simple(title, max_width_in, font_pt, bold=False, max_lines=2)
        wrapped_company = _wrap_text_simple(company, max_width_in, font_pt, bold=False, max_lines=1)
        return wrapped_title + "\n" + wrapped_company

    only_text = title or company
    if not only_text:
        return ""
    return _wrap_text_simple(only_text, max_width_in, font_pt, bold=False, max_lines=2)


def build_pptx_slide(
    prs: Presentation,
    template: Image.Image,
    slot_shape: SlotShape,
    slots: List[Slot],
    speakers: list,   # list of dicts: name, title, company, role, photo (PIL Image, already bust-cropped)
    session_name: str = "",
    hall_name: str = "",
    date_str: str = "",
):
    """
    Adds ONE slide to `prs` for this session. Mirrors slide_compositor.py's
    compose_slide() logic (same layout slots, same font-size rules, same
    theme detection) but emits real PowerPoint shapes instead of pixels.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # --- Background: the designer's template PNG, stretched to fill the slide ---
    bg_stream = _save_pil_image_to_stream(template)
    slide.shapes.add_picture(bg_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)

    W_px, H_px = template.size
    px_to_in = SLIDE_WIDTH_IN / W_px  # uniform scale since background is stretched to fill

    canvas_top_in = H_px * CANVAS_TOP_RATIO * px_to_in

    title_theme = detect_theme(template, (0, int(H_px * CANVAS_TOP_RATIO), W_px, min(H_px, int(H_px * (CANVAS_TOP_RATIO + 0.12)))))
    caption_theme = detect_theme(template, (0, int(H_px * (CANVAS_TOP_RATIO + 0.15)), W_px, H_px))

    # --- Session title ---
    title_reserved_in = 0.0
    if session_name:
        max_title_width_in = SLIDE_WIDTH_IN * 0.94
        title_pt = _fit_font_pt(session_name, max_title_width_in, SESSION_NAME_IDEAL_PT, SESSION_NAME_MIN_PT, bold=True)
        title_height_in = 0.40
        _add_text_box(
            slide, SLIDE_WIDTH_IN * 0.03, canvas_top_in, max_title_width_in, title_height_in,
            session_name, title_pt, bold=True, color_tuple=title_theme.primary_color,
            align=PP_ALIGN.LEFT,
        )
        # Tight, near single-line-spacing gap before the hall/date line -
        # was 0.55in (way too much vertical air), now just enough to clear
        # the title's own descenders.
        title_reserved_in = title_height_in + 0.03

    meta_parts = [p for p in [hall_name, date_str] if p]
    if meta_parts:
        meta_text = "   |   ".join(meta_parts)
        _add_text_box(
            slide, SLIDE_WIDTH_IN * 0.03, canvas_top_in + title_reserved_in,
            SLIDE_WIDTH_IN * 0.9, 0.3, meta_text, META_LINE_PT, bold=False,
            color_tuple=title_theme.secondary_color, align=PP_ALIGN.LEFT,
        )
        title_reserved_in += 0.32

    speaker_area_top_in = canvas_top_in + title_reserved_in
    speaker_area_height_in = SLIDE_HEIGHT_IN - speaker_area_top_in

    # Compute ONE shared font size for ALL speaker names up front (point 1
    # from project feedback) - rather than each name fitting itself
    # independently, every name starts at the same size, sized to fit
    # whichever name is longest/hardest to fit. This makes "select all name
    # boxes, change the font size once" behave the way the team expects.
    name_caption_width_in = (slots[0].w * SLIDE_WIDTH_IN * 0.96) if slots else SLIDE_WIDTH_IN
    shared_name_pt = _shared_name_font_pt(
        [sp["name"] for sp in speakers], name_caption_width_in,
        SPEAKER_NAME_PT, SPEAKER_NAME_PT * 0.6,
    )

    # --- Speaker slots ---
    for slot, sp in zip(slots, speakers):
        slot_x_in = slot.x * SLIDE_WIDTH_IN
        slot_y_in = speaker_area_top_in + slot.y * speaker_area_height_in
        slot_w_in = slot.w * SLIDE_WIDTH_IN
        slot_h_in = slot.h * speaker_area_height_in

        role_text = sp.get("role", "")
        show_role = role_text and role_text.strip().lower() not in DEFAULT_LABEL_VALUES_TO_HIDE

        role_reserved_in = 0.0
        if show_role:
            role_height_in = 0.25
            _add_text_box(
                slide, slot_x_in, slot_y_in, slot_w_in, role_height_in,
                role_text, ROLE_LABEL_PT, bold=True, color_tuple=caption_theme.role_label_color,
                align=PP_ALIGN.CENTER,
            )
            role_reserved_in = role_height_in

        caption_reserve_in = slot_h_in * 0.34
        photo_h_in = slot_h_in - role_reserved_in - caption_reserve_in
        photo_top_in = slot_y_in + role_reserved_in

        if sp.get("photo") is not None:
            # Apply the exact designer mask shape via Pillow (flattened,
            # per project decision), THEN insert as a PowerPoint picture.
            photo_w_px = max(1, int(slot_w_in / px_to_in))
            photo_h_px = max(1, int(photo_h_in / px_to_in))
            masked = _apply_mask_shape(sp["photo"], slot_shape, photo_w_px, photo_h_px)
            photo_stream = _save_pil_image_to_stream(masked, preserve_alpha=True)
            slide.shapes.add_picture(
                photo_stream, Inches(slot_x_in), Inches(photo_top_in),
                width=Inches(slot_w_in), height=Inches(photo_h_in),
            )

        caption_y_in = photo_top_in + photo_h_in + 0.04
        caption_max_width_in = slot_w_in * 0.96

        name_height_in = 0.3
        _add_text_box(
            slide, slot_x_in, caption_y_in, slot_w_in, name_height_in,
            sp["name"], shared_name_pt, bold=True, color_tuple=caption_theme.name_accent_color,
            align=PP_ALIGN.CENTER,
        )
        # Tighter name -> title/company gap (point 5 from project feedback;
        # was effectively ~0.0in already via no explicit gap, but the name
        # box's own height (0.3in) was adding visible whitespace below
        # short text - shrinking the box height directly closes that gap).
        caption_y_in += name_height_in - 0.06

        title_company_text = _build_title_company_lines(
            sp.get("title", ""), sp.get("company", ""), caption_max_width_in, TITLE_COMPANY_PT,
        )
        if title_company_text:
            line_count = title_company_text.count("\n") + 1
            box_height_in = 0.22 * line_count
            _add_multiline_text_box(
                slide, slot_x_in, caption_y_in, slot_w_in, box_height_in,
                title_company_text, TITLE_COMPANY_PT, bold=False,
                color_tuple=caption_theme.caption_color, align=PP_ALIGN.CENTER,
            )

    return slide


def save_pptx(prs: Presentation, output_path: str):
    prs.save(output_path)
