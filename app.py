from pathlib import Path
base = Path('output/soft_slides_app')
app_code = '''import io, math, re
from pathlib import Path
import streamlit as st
import pandas as pd
from PIL import Image, ImageOps, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches

BASE_DIR = Path(__file__).resolve().parent
OUTPUT_DIR = BASE_DIR / 'output'
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

st.set_page_config(page_title='Soft Slides Generator', layout='wide')
st.title('Conference Soft Slides Generator')
st.caption('Upload template, speaker cutout, Excel, photos, and optional font.')

TEMPLATE = st.file_uploader('Upload flat template image', type=['png', 'jpg', 'jpeg'])
CUTOUT = st.file_uploader('Upload single speaker cutout PNG (transparent)', type=['png'])
EXCEL = st.file_uploader('Upload speaker Excel', type=['xlsx', 'xls'])
PHOTOS = st.file_uploader('Upload speaker photos folder (multiple files)', type=['png', 'jpg', 'jpeg'], accept_multiple_files=True)
FONT_FILE = st.file_uploader('Upload font file (optional)', type=['ttf', 'otf'])
EVENT_NAME = st.text_input('Event / session title', value='')
HALL_NAME = st.text_input('Hall name', value='')
DATE_TEXT = st.text_input('Date text', value='')
SHOW_LEADERSHIP = st.checkbox('Show leadership roles above photos only', value=True)

if 'order' not in st.session_state:
    st.session_state.order = []
if 'version' not in st.session_state:
    st.session_state.version = 1
if 'photo_map_manual' not in st.session_state:
    st.session_state.photo_map_manual = {}

ROLE_PRIORITY = {'MODERATOR': 0, 'CO-MODERATOR': 1, 'CHAIR': 2, 'CO-CHAIR': 3, 'KEYNOTE': 4, 'KEYNOTE SPEAKER': 4, 'PANELIST': 5, 'SPEAKER': 6}
LEADERSHIP_ROLES = {'MODERATOR', 'CO-MODERATOR', 'CHAIR', 'CO-CHAIR', 'KEYNOTE', 'KEYNOTE SPEAKER'}
FIELD_HINT = {'role': ['role'], 'speaker name': ['speaker name', 'name'], 'title': ['title', 'designation'], 'company': ['company', 'org', 'organization']}


def safe(v):
    return '' if pd.isna(v) else str(v)


def clean(s):
    s = re.sub(r'[^a-z0-9]+', ' ', safe(s).lower())
    return re.sub(r'\s+', ' ', s).strip()


def tokens(s):
    return [t for t in clean(s).split() if t]


def pick_col(cols, candidates):
    lower = {str(c).lower().strip(): c for c in cols}
    for cand in candidates:
        if cand in lower:
            return lower[cand]
    return cols[0] if len(cols) else None


def role_rank(role):
    r = clean(role).upper()
    for k, v in ROLE_PRIORITY.items():
        kk = clean(k).upper()
        if r == kk or kk in r:
            return v
    return 99


def ordered_indices_by_role(df):
    role_col = pick_col(df.columns, FIELD_HINT['role'])
    if not role_col:
        return list(df.index)
    roles = df[role_col].astype(str).fillna('')
    return sorted(df.index.tolist(), key=lambda i: (role_rank(roles.loc[i]), i))


def text_size(draw, txt, font):
    box = draw.textbbox((0, 0), txt, font=font)
    return box[2] - box[0], box[3] - box[1]


def get_font(font_path, size):
    if font_path and font_path.exists():
        try:
            return ImageFont.truetype(str(font_path), int(round(size)))
        except Exception:
            pass
    for p in ['/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf', '/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf']:
        try:
            return ImageFont.truetype(p, int(round(size)))
        except Exception:
            pass
    return ImageFont.load_default()


def fit_font(draw, texts, font_path, start=12, min_size=7.5, max_width=300):
    size = start
    while size >= min_size:
        f = get_font(font_path, size)
        if all(text_size(draw, safe(t), f)[0] <= max_width for t in texts):
            return size
        size -= 0.5
    return min_size


def display_name(v):
    v = safe(v).strip()
    parts = v.split()
    if len(parts) >= 3 and len(v) > 20:
        return parts[0] + ' ' + parts[1][0] + '.'
    if len(v) > 28:
        return v[:25].rstrip() + '...'
    return v


def parse_row(row):
    vals = [safe(x).strip() for x in row.tolist() if safe(x).strip()]
    role = ''
    name = ''
    extras = []
    for v in vals:
        up = v.upper()
        if any(k in up for k in ['MODERATOR', 'CO-MODERATOR', 'CHAIR', 'CO-CHAIR', 'KEYNOTE', 'PANELIST', 'SPEAKER']):
            role = v
        elif not name:
            name = v
        else:
            extras.append(v)
    title, company = '', ''
    if len(extras) >= 2:
        title, company = extras[0], extras[1]
    elif len(extras) == 1:
        tc = extras[0]
        if '|' in tc:
            title, company = [x.strip() for x in tc.split('|', 1)]
        elif ',' in tc:
            title, company = [x.strip() for x in tc.split(',', 1)]
        else:
            title = tc
    return role, name, title, company


def normalize_photo_key(name):
    s = clean(name)
    s = s.replace('h e ', '').replace('he ', '')
    s = s.replace('shri ', '').replace('smt ', '').replace('dr ', '').replace('prof ', '')
    s = s.replace('mr ', '').replace('mrs ', '').replace('ms ', '')
    return clean(s)


def score_match(name, filename):
    n = normalize_photo_key(name)
    f = normalize_photo_key(filename)
    nt = tokens(n)
    ft = tokens(f)
    if not n or not f:
        return 0
    if n == f:
        return 100
    if nt and ft and nt[-1] == ft[-1] and len(nt[-1]) > 2:
        return 90
    common = len(set(nt) & set(ft))
    if common:
        return 20 + common * 10
    if n in f or f in n:
        return 50
    return 0


def assign_photos(df, order, photo_map):
    assigned = {}
    used = set()
    rows = [parse_row(df.iloc[i]) for i in order]
    for idx, row in zip(order, rows):
        name = row[1]
        best = None
        best_score = 0
        for fname, fileobj in photo_map.items():
            if fname in used:
                continue
            sc = score_match(name, fname)
            if sc > best_score:
                best = fname
                best_score = sc
        if best_score >= 20:
            assigned[idx] = best
            used.add(best)
        else:
            assigned[idx] = None
    return assigned, rows


def role_label(role):
    r = clean(role).upper()
    if 'CO-MODERATOR' in r:
        return 'CO-MODERATOR'
    if 'MODERATOR' in r:
        return 'MODERATOR'
    if 'CO-CHAIR' in r:
        return 'CO-CHAIR'
    if 'CHAIR' in r:
        return 'CHAIR'
    if 'KEYNOTE' in r:
        return 'KEYNOTE SPEAKER'
    return ''


def group_label(rows):
    roles = [r[0] for r in rows if r[0]]
    if any(clean(r).upper().startswith('PANELIST') for r in roles):
        return 'PANELISTS'
    if any(clean(r).upper().startswith('SPEAKER') for r in roles):
        return 'SPEAKERS'
    return 'SPEAKERS'


def layout(n):
    if n <= 2:
        return 2, 1
    if n <= 4:
        return 2, 2
    if n <= 6:
        return 3, 2
    if n <= 9:
        return 3, 3
    return 4, 4


if EXCEL is not None:
    df = pd.read_excel(EXCEL)
    st.subheader('Speaker Data')
    st.dataframe(df, use_container_width=True)

    if st.button('Auto order by role'):
        st.session_state.order = ordered_indices_by_role(df)
    if not st.session_state.order:
        st.session_state.order = list(df.index)

    st.subheader('Reorder speakers')
    st.caption('Leadership roles stay first when possible.')
    for i, idx in enumerate(list(st.session_state.order)):
        cols = st.columns([6, 1, 1])
        row = df.iloc[idx]
        row_role, row_name, _, _ = parse_row(row)
        cols[0].write(f'{i+1}. {row_name or safe(row.iloc[0])} {"(" + row_role + ")" if row_role else ""}')
        if cols[1].button('▲', key=f'u{idx}', disabled=i == 0):
            st.session_state.order[i - 1], st.session_state.order[i] = st.session_state.order[i], st.session_state.order[i - 1]
            st.rerun()
        if cols[2].button('▼', key=f'd{idx}', disabled=i == len(st.session_state.order) - 1):
            st.session_state.order[i + 1], st.session_state.order[i] = st.session_state.order[i], st.session_state.order[i + 1]
            st.rerun()

    if PHOTOS:
        st.subheader('Photo matching review')
        photo_map = {Path(f.name).stem: f for f in PHOTOS}
        assigned, parsed_rows = assign_photos(df, st.session_state.order, photo_map)
        photo_names = list(photo_map.keys())
        for idx, row in zip(st.session_state.order, parsed_rows):
            role, name, title, company = row
            matched = assigned[idx]
            cols = st.columns([3, 2, 4])
            cols[0].write(name)
            cols[1].write(matched if matched else 'NO MATCH')
            cols[2].write(role_label(role) or '')

    if st.button('Generate Downloads'):
        if TEMPLATE is None or CUTOUT is None:
            st.error('Template and cutout PNG are required.')
            st.stop()
        if not PHOTOS:
            st.error('Please upload speaker photos.')
            st.stop()

        photo_map = {Path(f.name).stem: f for f in PHOTOS}
        assigned, parsed = assign_photos(df, st.session_state.order, photo_map)
        base_img = Image.open(TEMPLATE).convert('RGBA')
        cutout_img = Image.open(CUTOUT).convert('RGBA')
        cutout_size = cutout_img.size
        W, H = base_img.size
        base_draw = ImageDraw.Draw(base_img)

        font_path = None
        if FONT_FILE is not None:
            font_path = OUTPUT_DIR / FONT_FILE.name
            font_path.write_bytes(FONT_FILE.getbuffer())

        names = [p[1] for p in parsed]
        name_size = fit_font(base_draw, names, font_path, start=14, min_size=8, max_width=int(W * 0.16))
        title_size = max(8, name_size - 2)
        role_size = max(8, title_size)

        def render_slide(indices, parsed_rows):
            img = base_img.copy()
            d = ImageDraw.Draw(img)
            for txt, base_fs, y_ratio in [(EVENT_NAME, 24, 0.07), (HALL_NAME, 18, 0.13), (DATE_TEXT, 16, 0.17)]:
                if txt.strip():
                    fs = fit_font(d, [txt], font_path, start=base_fs, min_size=8, max_width=int(W * 0.75))
                    f = get_font(font_path, fs)
                    tw, _ = text_size(d, txt, f)
                    d.text(((W - tw) / 2, int(H * y_ratio)), txt, font=f, fill=(255, 255, 255, 255))

            n = len(indices)
            cols, rows = layout(n)
            x_margin = int(W * 0.05)
            y_start = int(H * 0.24)
            usable_w = W - 2 * x_margin
            usable_h = int(H * 0.62)
            cell_w = usable_w / cols
            cell_h = usable_h / rows
            photo_w = int(min(cell_w, cell_h) * 0.68)
            photo_h = int(photo_w * cutout_size[1] / cutout_size[0]) if cutout_size[0] else photo_w
            photo_h = max(photo_h, 1)

            group = group_label(parsed_rows)
            d.text((int(W * 0.50), int(H * 0.22)), group, font=get_font(font_path, 24), fill=(255, 255, 255, 255), anchor='mm')

            for pos, (idx, rowdata) in enumerate(zip(indices, parsed_rows)):
                role, name, title, company = rowdata
                row, col = divmod(pos, cols)
                x = int(x_margin + col * cell_w + (cell_w - photo_w) / 2)
                y = int(y_start + row * cell_h)
                matched_key = assigned.get(idx)
                if matched_key:
                    ph = Image.open(photo_map[matched_key]).convert('RGBA')
                    ph = ImageOps.fit(ph, (photo_w, photo_h), method=Image.Resampling.LANCZOS, centering=(0.5, 0.33))
                    mask = Image.open(CUTOUT).convert('RGBA').resize((photo_w, photo_h))
                    alpha = mask.split()[-1] if mask.mode == 'RGBA' else None
                    if alpha is None:
                        alpha = Image.new('L', (photo_w, photo_h), 255)
                    img.paste(ph, (x, y), alpha)
                else:
                    outline = Image.open(CUTOUT).convert('RGBA').resize((photo_w, photo_h))
                    img.alpha_composite(outline, (x, y))
                    d.text((x, y + photo_h + 2), f'No photo for {name}', font=get_font(font_path, 8), fill=(255, 180, 180, 255))

                nf = get_font(font_path, name_size)
                tf = get_font(font_path, title_size)
                cf = get_font(font_path, title_size)
                rf = get_font(font_path, role_size)
                disp = display_name(name)
                nw, _ = text_size(d, disp, nf)
                d.text((x + (photo_w - nw) / 2, y + photo_h + 8), disp, font=nf, fill=(255, 235, 80, 255))
                rl = role_label(role)
                if rl in LEADERSHIP_ROLES:
                    d.text((x + photo_w / 2, y - 22), rl, font=rf, fill=(255, 255, 255, 255), anchor='mm')
                d.text((x, y + photo_h + 25), title, font=tf, fill=(240, 240, 240, 255))
                d.text((x, y + photo_h + 43), company, font=cf, fill=(240, 240, 240, 255))
            return img

        slide_img = render_slide(st.session_state.order, parsed)
        buf = io.BytesIO()
        slide_img.save(buf, format='PNG')

        ppt = Presentation()
        ppt.slide_width = Inches(13.333)
        ppt.slide_height = Inches(7.5)
        slide = ppt.slides.add_slide(ppt.slide_layouts[6])
        tmp = OUTPUT_DIR / f'softslide_{st.session_state.version}_1.png'
        tmp.write_bytes(buf.getvalue())
        slide.shapes.add_picture(str(tmp), 0, 0, width=ppt.slide_width, height=ppt.slide_height)

        ppt_buf = io.BytesIO()
        ppt.save(ppt_buf)
        ppt_buf.seek(0)
        st.download_button('Download PNG', data=buf.getvalue(), file_name=f'soft_slide_v{st.session_state.version}.png', mime='image/png')
        st.download_button('Download PPTX', data=ppt_buf.getvalue(), file_name=f'soft_slide_v{st.session_state.version}.pptx', mime='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        st.session_state.version += 1
        st.success('Generated 1 slide.')
'''
(base / 'app.py').write_text(app_code)
(base / 'script.py').write_text('from pathlib import Path\nimport runpy\nrunpy.run_path(str(Path(__file__).with_name("app.py")), run_name="__main__")\n')
(base / 'requirements.txt').write_text('streamlit\npandas\npillow\npython-pptx\nopenpyxl\n')
print('written')
